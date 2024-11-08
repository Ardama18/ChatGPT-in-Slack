from PyMuPDF import fitz  # PyMuPDF
from pathlib import Path
from docx import Document
from pptx import Presentation
from typing import TypedDict

class Doc(TypedDict):
  stem: str
  body: str


def read_text_as_utf8_or_sjis(text_path):
  """utf-8で読んでだめならshift_jisで読む"""
  try:
    body = text_path.read_text()
  except UnicodeDecodeError:
    body = text_path.read_text(encoding="shift_jis")
  return body



def load_document(document_dir=Path("data/")) -> list[Doc]:
  document_list = []
  for document_path in document_dir.glob("*"):
    body: str = ""
    if ".docx" in document_path.suffixes:
      document = Document(str(document_path))
      for _, p in enumerate(document.paragraphs):
        if p.text != "":
          body += p.text
    elif ".pptx" in document_path.suffixes:
      pptx = Presentation(str(document_path))
      for _, slide in enumerate(pptx.slides):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
          textNote = slide.notes_slide.notes_text_frame.text
          if textNote != "":
            body += textNote
    elif ".txt" in document_path.suffixes:
      try:
        body = document_path.read_text()
      except UnicodeDecodeError:
        body = document_path.read_text(encoding="shift_jis")
    elif ".pdf" in document_path.suffixes:
        pdf_document = fitz.open(str(document_path))
        for page_num in range(pdf_document.page_count):
            page = pdf_document.load_page(page_num)
            body += page.get_text()  
    else:
      raise RuntimeError("想定外のファイル")

    one_document: Doc = {
      "stem": document_path.stem,
      "body": body
    }
    document_list.append(one_document)
  return document_list


def concat_document(document_dir=Path("data/")):
  """ドキュメントをコンテキストウィンドウに乗っけるため、XML形式でまとめる"""

  docs = load_document(document_dir)
  document_list: list[str] = []
  for d in docs:
    stem = d["stem"]
    body = d["body"]
    one_document_text = f"<document><meta>{stem}</meta><body>{body}</body></document>"
    document_list.append(one_document_text)
  return "\n".join(document_list)
